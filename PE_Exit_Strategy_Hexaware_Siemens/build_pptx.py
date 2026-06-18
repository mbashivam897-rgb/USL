"""
Builds: PE_Exit_Strategy_Hexaware_Siemens.pptx
Institutional PE Exit Strategy deck (16:9), themed, ~26 slides.
Figures are consistent with PE_Exit_Model_Hexaware_Siemens.xlsx.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

NAVY  = RGBColor(0x1F, 0x38, 0x64)
BLUE  = RGBColor(0x2E, 0x54, 0x96)
LBLUE = RGBColor(0xD6, 0xE0, 0xF0)
GOLD  = RGBColor(0xBF, 0x90, 0x00)
LGOLD = RGBColor(0xFF, 0xF2, 0xCC)
GREY  = RGBColor(0xF2, 0xF2, 0xF2)
DGREY = RGBColor(0x59, 0x59, 0x59)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GREEN = RGBColor(0x2E, 0x7D, 0x32)
RED   = RGBColor(0xC0, 0x00, 0x00)
BLACK = RGBColor(0x00, 0x00, 0x00)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def add_slide():
    return prs.slides.add_slide(BLANK)

def rect(slide, l, t, w, h, color, line=None):
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line; sp.line.width = Pt(0.75)
    sp.shadow.inherit = False
    return sp

def txt(slide, l, t, w, h, text, size=14, color=BLACK, bold=False, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, italic=False, font="Calibri"):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
    f.color.rgb = color; f.name = font
    return tb

def bullets(slide, l, t, w, h, items, size=14, color=BLACK, gap=4):
    tb = slide.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True
    first = True
    for it in items:
        lvl = 0; mark = "•  "
        if isinstance(it, tuple): it, lvl = it
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.level = lvl; p.space_after = Pt(gap)
        if lvl > 0: mark = "–  "
        r = p.add_run(); r.text = mark + it
        r.font.size = Pt(size - lvl); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def header(slide, kicker, title, accent=GOLD):
    rect(slide, 0, 0, SW, Inches(1.15), NAVY)
    rect(slide, 0, Inches(1.15), SW, Inches(0.06), accent)
    txt(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.35), kicker, 12, accent, bold=True)
    txt(slide, Inches(0.5), Inches(0.42), Inches(12.3), Inches(0.7), title, 24, WHITE, bold=True)
    # footer
    txt(slide, Inches(0.5), Inches(7.05), Inches(9), Inches(0.35),
        "Strictly Private & Confidential — IC Memorandum / Exit Advisory", 9, DGREY)
    txt(slide, Inches(11.0), Inches(7.05), Inches(2.0), Inches(0.35),
        f"{len(prs.slides.__iter__.__self__._sldIdLst)}", 9, DGREY, align=PP_ALIGN.RIGHT)

def table(slide, l, t, w, rows, col_widths=None, header_fill=BLUE, font=10.5,
          first_col_label=True, zebra=True, h=None):
    nrows = len(rows); ncols = len(rows[0])
    height = h or Inches(0.4 * nrows)
    gt = slide.shapes.add_table(nrows, ncols, l, t, w, height).table
    if col_widths:
        for i, cw in enumerate(col_widths):
            gt.columns[i].width = cw
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = gt.cell(ri, ci)
            cell.margin_left = Inches(0.06); cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.02); cell.margin_bottom = Inches(0.02)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.CENTER
            r = p.add_run(); r.text = str(val); r.font.name = "Calibri"
            if ri == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = header_fill
                r.font.bold = True; r.font.color.rgb = WHITE; r.font.size = Pt(font)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE if (not zebra or ri % 2 == 1) else GREY
                r.font.size = Pt(font); r.font.color.rgb = BLACK
                if ci == 0 and first_col_label:
                    r.font.bold = True; r.font.color.rgb = NAVY
    return gt

# =====================================================================
# 1. TITLE
# =====================================================================
s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, 0, Inches(2.6), SW, Inches(0.08), GOLD)
txt(s, Inches(0.8), Inches(0.7), Inches(12), Inches(0.5),
    "PRIVATE EQUITY — VALUE CREATION & EXIT ADVISORY", 14, GOLD, bold=True)
txt(s, Inches(0.8), Inches(2.8), Inches(11.8), Inches(1.6),
    "Exit Strategy Report", 46, WHITE, bold=True)
txt(s, Inches(0.8), Inches(4.0), Inches(11.8), Inches(1.0),
    "Hexaware Technologies Ltd.  |  Siemens Ltd. (India)", 26, LBLUE, bold=True)
txt(s, Inches(0.8), Inches(5.4), Inches(11.8), Inches(0.5),
    "IC Investment Memo  •  Sell-Side Exit Advisory  •  M&A Strategic Review", 15, WHITE)
txt(s, Inches(0.8), Inches(6.6), Inches(11.8), Inches(0.4),
    "Strictly Private & Confidential   •   18 June 2026", 12, GOLD)

# =====================================================================
# 2. METHODOLOGY & DISCLAIMER
# =====================================================================
s = add_slide()
header(s, "IMPORTANT", "Methodology, Sources & Key Assumptions")
bullets(s, Inches(0.5), Inches(1.45), Inches(12.4), Inches(3.3), [
    "Two distinct assets: (1) Hexaware — a LIVE PE case (Carlyle ~74% post-IPO, acquired 2021, partial IPO Feb-2025); (2) Siemens India — a STRATEGIC subsidiary (Siemens AG ~69%), analysed as a hypothetical PE / take-private study.",
    "Hexaware reports on a calendar year (Dec-end); Siemens on a fiscal year (Oct–Sep). Siemens FY20–24 figures include the now-demerged Energy business and are not directly comparable to continuing operations.",
    "FX assumption: INR/USD approx. 85 for FY25–26 conversions.",
    "Sources: company IR & press releases, NSE/BSE filings, Carlyle disclosures, NASSCOM, Kroll, Mordor Intelligence, Jefferies/Nuvama, WSJ, Stockopedia (all cited on the data slides and in the accompanying report/Excel model).",
    "All valuation outputs (EV, equity, MOIC, IRR) are analytical estimates built on stated assumptions and computed live in the companion Excel model. Gross of fees/carry. Not investment advice.",
], size=14)
rect(s, Inches(0.5), Inches(5.0), Inches(12.33), Inches(1.5), LGOLD)
txt(s, Inches(0.7), Inches(5.15), Inches(12), Inches(1.2),
    "Companion deliverables:  (1) this deck;  (2) a 14–15 page written report (Markdown);  "
    "(3) a formula-linked Excel model 'PE_Exit_Model_Hexaware_Siemens.xlsx' where Assumptions feed Valuation & Returns sheets and a live IC Dashboard.",
    13, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# 3. AGENDA
# =====================================================================
s = add_slide()
header(s, "CONTENTS", "Report Structure — 10 Sections per Company")
left_items = ["1.  Investment Overview","2.  Business & Industry Analysis","3.  Financial Performance (5-yr)",
              "4.  Value Creation Opportunities","5.  Exit Readiness Assessment"]
right_items = ["6.  Buyer Universe Analysis","7.  Exit Strategy Analysis","8.  Exit Valuation (Bear/Base/Bull)",
               "9.  Exit Timing Analysis","10. Final Recommendation"]
bullets(s, Inches(0.8), Inches(1.7), Inches(5.8), Inches(4.5), left_items, size=18, gap=14, color=NAVY)
bullets(s, Inches(6.9), Inches(1.7), Inches(5.8), Inches(4.5), right_items, size=18, gap=14, color=NAVY)
rect(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.8), LBLUE)
txt(s, Inches(1.0), Inches(6.2), Inches(11.3), Inches(0.6),
    "Plus: Porter's Five Forces & SWOT (Sec 2), comparative scorecard, and consolidated IC recommendation.",
    13, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# =====================================================================
# 4. EXECUTIVE SUMMARY
# =====================================================================
s = add_slide()
header(s, "EXECUTIVE SUMMARY", "Two Assets, Two Very Different PE Verdicts")
table(s, Inches(0.5), Inches(1.5), Inches(12.33), [
    ["Dimension","Hexaware Technologies","Siemens Ltd. (India)"],
    ["Sector","IT / digital services","Industrial tech / electrification"],
    ["PE relevance","LIVE case — Carlyle ~74%","Strategic parent — Siemens AG ~69%"],
    ["Control available to PE?","YES (sponsor holds it)","NO"],
    ["Exit readiness (/10)","8.4","9.0 (but no sponsor control)"],
    ["Valuation backdrop","Near-trough (AI de-rating)","Premium (richly valued)"],
    ["Base equity value","~Rs 46,000 cr (~US$5.4bn)","~Rs 166,500 cr (~US$19.6bn)"],
    ["Base MOIC / IRR","~2.0x / ~11%","~1.28x / ~6.4% (hypothetical)"],
    ["IC VERDICT","HOLD & STAGGER-EXIT by ~CY2028","PASS on control"],
], col_widths=[Inches(3.0), Inches(4.66), Inches(4.66)], font=12.5, h=Inches(4.7))
txt(s, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.6),
    "Deploy effort against Hexaware-type sponsor-controlled, re-rating-optional assets; study Siemens for the value-unlock (demerger) template, but do not chase a control position you cannot obtain.",
    12, DGREY, italic=True)

# =====================================================================
# SECTION DIVIDER helper
# =====================================================================
def divider(part, title, subtitle):
    sl = add_slide()
    rect(sl, 0, 0, SW, SH, NAVY)
    rect(sl, Inches(0.8), Inches(3.0), Inches(2.2), Inches(0.10), GOLD)
    txt(sl, Inches(0.8), Inches(2.0), Inches(11), Inches(0.6), part, 16, GOLD, bold=True)
    txt(sl, Inches(0.8), Inches(3.2), Inches(11.5), Inches(1.2), title, 40, WHITE, bold=True)
    txt(sl, Inches(0.8), Inches(4.6), Inches(11.5), Inches(0.8), subtitle, 18, LBLUE)
    return sl

# =====================================================================
# PART A — HEXAWARE
# =====================================================================
divider("PART A", "Hexaware Technologies Ltd.",
        "NSE: HEXT  •  Sponsor: The Carlyle Group (via CA Magnum Holdings)")

# A.1 Investment Overview
s = add_slide()
header(s, "HEXAWARE — SECTION 1", "Investment Overview & the PE Story")
bullets(s, Inches(0.5), Inches(1.4), Inches(6.1), Inches(5.2), [
    "Global AI-first IT & BPS provider: 31,000+ staff, ~54 offices, 28 countries.",
    "Proprietary platforms: RapidX (digital), Tensai (AI/automation), Amaze (cloud).",
    "Verticals: BFSI, Healthcare/Insurance, Manufacturing, Hi-Tech, Travel/Logistics.",
    "Geography: North America ~70%+; Europe & APAC the expansion runway.",
    "Strengths: scaled growth, AI brand, blue-chip logos, net-cash, proven mgmt.",
    "Risks: AI cannibalisation of headcount model, growth deceleration, sponsor overhang, FX.",
], size=13.5)
rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.2), GREY)
txt(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.4), "PE OWNERSHIP TIMELINE", 13, NAVY, bold=True)
bullets(s, Inches(7.0), Inches(1.95), Inches(5.7), Inches(4.5), [
    "2013: Baring PE Asia acquires ~70.8% control.",
    "2020: Delisted (BPEA take-private).",
    "Oct-2021: Carlyle buys at ~US$3bn EV (~Rs 23,200 cr).",
    "Feb-2025: Rs 8,750 cr (~US$1bn) IPO, 100% OFS; valued ~Rs 43,000 cr (~US$5.0–5.4bn).",
    "Largest tech-services IPO globally in over a decade; largest sponsor-owned IPO in India.",
    "Carlyle stake: 95% -> ~74.1% post-IPO — bulk of value still un-monetised.",
], size=12.5)

# A.2 Business & Industry + Porter's + SWOT
s = add_slide()
header(s, "HEXAWARE — SECTION 2", "Industry, Porter's Five Forces & SWOT")
txt(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
    "India tech sector FY25 ~US$282.6bn (+5.1%), >US$300bn FY26 (NASSCOM); India IT services US$37bn (2025) -> US$57bn (2030), ~9% CAGR (Mordor).",
    11.5, DGREY, italic=True)
table(s, Inches(0.5), Inches(1.8), Inches(6.0), [
    ["Porter's Force","Intensity"],
    ["Competitive rivalry","High"],
    ["New entrants","Low–Medium"],
    ["Supplier (talent) power","Medium–High"],
    ["Buyer power","High"],
    ["Substitutes (GenAI)","Med–High & rising"],
], col_widths=[Inches(4.0), Inches(2.0)], font=12, h=Inches(3.0))
txt(s, Inches(6.8), Inches(1.8), Inches(6), Inches(0.35), "SWOT", 14, NAVY, bold=True)
table(s, Inches(6.8), Inches(2.15), Inches(6.0), [
    ["Strengths","Weaknesses"],
    ["Scaled growth, AI brand, platforms, net-cash","US/BFSI concentration, mid-tier margins, overhang"],
    ["Opportunities","Threats"],
    ["GenAI wave, vendor consolidation, Europe, M&A","AI cannibalisation, de-rating, macro, talent/FX"],
], col_widths=[Inches(3.0), Inches(3.0)], font=10.5, first_col_label=False, h=Inches(2.6))
txt(s, Inches(0.5), Inches(5.6), Inches(6.1), Inches(1.4),
    "Net: structurally competitive industry with a RISING AI-substitution threat — the central swing factor for exit multiple and timing.",
    12.5, NAVY, bold=True)

# A.3 Financials
s = add_slide()
header(s, "HEXAWARE — SECTION 3", "5-Year Financial Performance (CY, Dec-end)")
table(s, Inches(0.5), Inches(1.5), Inches(12.33), [
    ["Metric","CY21","CY22","CY23","CY24","CY25","CAGR"],
    ["Revenue (Rs cr)","7,178","9,200","10,380","11,974","13,430","~17%"],
    ["Revenue (US$ m)","970","1,156","1,257","1,429","1,537","~12%"],
    ["Rev growth % (USD)","—","~19%","~9%","+13.7%","+7.6%","—"],
    ["EBITDA margin %","~17%","~16.5%","~16.5%","~17%","~16.7%","stable"],
    ["EBITDA (Rs cr, est.)","~1,220","~1,500","~1,700","~2,050","~2,250","—"],
    ["PAT (Rs cr, est.)","750","820","998","1,150","1,300","~15%"],
], col_widths=[Inches(2.9)] + [Inches(1.57)]*6, font=11.5, h=Inches(3.0))
bullets(s, Inches(0.5), Inches(4.8), Inches(12.3), Inches(1.7), [
    "Revenue nearly DOUBLED CY21->CY25 — validates the PE value-creation thesis.",
    "But growth DECELERATED to +7.6% (USD) in CY25 — a yellow flag for the exit multiple.",
    "EBITDA ~16–17% (below tier-1's ~24–26%) = a margin-expansion lever; ROE ~22–25%, net-cash B/S = an LBO leverage lever.",
], size=13)
txt(s, Inches(0.5), Inches(6.7), Inches(12), Inches(0.4),
    "Sources: Stockopedia (INR revenue); Hexaware/PR Newswire (USD revenue); EBITDA & PAT analyst estimates.", 10, DGREY, italic=True)

# A.4 Value creation
s = add_slide()
header(s, "HEXAWARE — SECTION 4", "Value Creation Opportunities")
table(s, Inches(0.5), Inches(1.5), Inches(12.33), [
    ["Lever","Action","Impact"],
    ["Margin improvement","Pyramid/utilisation, automation, offshore mix","+150–300 bps EBITDA (toward ~19–20%)"],
    ["Revenue expansion","Large-deal TCV, GenAI services, Europe push","+200–400 bps revenue CAGR"],
    ["Digital / AI","Monetise Tensai/RapidX/Amaze IP; outcome pricing","Multiple re-rating ('AI-native')"],
    ["Strategic M&A","Bolt-ons in domain / AI / data","Inorganic growth + capability"],
    ["Capital structure","Introduce leverage in a secondary LBO recap","Equity-return amplification"],
], col_widths=[Inches(2.8), Inches(5.6), Inches(3.9)], font=11.5, h=Inches(3.0))
rect(s, Inches(0.5), Inches(5.0), Inches(12.33), Inches(1.5), LGOLD)
txt(s, Inches(0.7), Inches(5.15), Inches(12), Inches(1.2),
    "Illustrative: lifting EBITDA margin from ~17% to ~20% on ~Rs 15,000 cr revenue adds ~Rs 450 cr EBITDA; "
    "at 16–18x EV/EBITDA that is ~Rs 7,000–8,000 cr of incremental enterprise value — before any growth re-acceleration or re-rating.",
    13, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# A.5 Exit readiness
s = add_slide()
header(s, "HEXAWARE — SECTION 5", "Exit Readiness Assessment — 8.4 / 10")
table(s, Inches(0.5), Inches(1.5), Inches(7.4), [
    ["Dimension","Score /10"],
    ["Management quality","9"],["Corporate governance","9"],
    ["Scalability","8"],["Market positioning","7"],
    ["Investor appeal","7"],["Strategic attractiveness","8"],
    ["Financial maturity","9"],["Public-market readiness","10"],
    ["OVERALL","8.4"],
], col_widths=[Inches(5.4), Inches(2.0)], font=11.5, h=Inches(4.6))
rect(s, Inches(8.2), Inches(1.5), Inches(4.6), Inches(4.6), GREY)
txt(s, Inches(8.4), Inches(1.7), Inches(4.2), Inches(4.2),
    "Maximally exit-ready: the asset is ALREADY public, governed and liquid.\n\n"
    "The binding constraint is the market MULTIPLE, not readiness.\n\n"
    "=> The live questions are exit ROUTE and TIMING, not whether the asset can be exited.",
    15, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# A.6 Buyer universe
s = add_slide()
header(s, "HEXAWARE — SECTION 6", "Buyer Universe — Ranked by Probability")
txt(s, Inches(0.5), Inches(1.3), Inches(6), Inches(0.4), "Strategic buyers", 14, NAVY, bold=True)
table(s, Inches(0.5), Inches(1.7), Inches(6.0), [
    ["Buyer","Likelihood"],
    ["Cognizant","Medium-High"],["Capgemini","Medium"],
    ["Accenture","Low-Medium"],["IBM / DXC","Low-Medium"],
], col_widths=[Inches(4.0), Inches(2.0)], font=11.5, h=Inches(2.2))
txt(s, Inches(6.8), Inches(1.3), Inches(6), Inches(0.4), "Financial buyers", 14, NAVY, bold=True)
table(s, Inches(6.8), Inches(1.7), Inches(6.0), [
    ["Buyer","Likelihood"],
    ["Blackstone","High"],["EQT (ex-BPEA)","Medium-High"],
    ["KKR / Bain / Advent","Medium"],["Brookfield / SWFs","High (blocks)"],
], col_widths=[Inches(4.0), Inches(2.0)], font=11.5, h=Inches(2.2))
rect(s, Inches(0.5), Inches(4.4), Inches(12.33), Inches(2.0), LBLUE)
bullets(s, Inches(0.7), Inches(4.55), Inches(12), Inches(1.8), [
    "Ranked probability: (1) capital-markets block sale to FIIs/DIIs/SWFs; (2) sponsor-to-sponsor secondary; (3) strategic sale to Cognizant/Capgemini; (4) IBM/Accenture tuck-in; (5) DXC/smaller strategics.",
    "A full strategic acquisition of a ~US$5bn LISTED company needs an open-offer/control premium and is operationally complex — more likely a longer-horizon scenario.",
], size=13, color=NAVY)

# A.7 Exit routes
s = add_slide()
header(s, "HEXAWARE — SECTION 7", "Exit Strategy Analysis — Routes Ranked")
table(s, Inches(0.5), Inches(1.5), Inches(12.33), [
    ["Route","Complexity","Valuation impact","Feasibility","P(success)"],
    ["1. Staggered block sales (listed mkt)","Low-Med","Neutral","High","~85%"],
    ["2. Secondary / sponsor-to-sponsor","Medium","Market + modest premium","Medium-High","~55%"],
    ["3. Strategic sale","High","+ premium","Medium","~25%"],
    ["4. Partial exit","Low","Neutral","High","~85%"],
    ["5. Merger","Very High","Variable","Low","~10%"],
], col_widths=[Inches(4.3), Inches(1.8), Inches(2.9), Inches(1.83), Inches(1.5)], font=11, h=Inches(3.0))
rect(s, Inches(0.5), Inches(5.0), Inches(12.33), Inches(1.4), LGOLD)
txt(s, Inches(0.7), Inches(5.15), Inches(12), Inches(1.1),
    "Recommended: STAGGERED secondary block sales via the listed market (primary), with a sponsor-to-sponsor full block as the opportunistic alternative and strategic sale retained as premium optionality.",
    14, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# A.8 Valuation
s = add_slide()
header(s, "HEXAWARE — SECTION 8", "Exit Valuation — Bear / Base / Bull")
table(s, Inches(0.5), Inches(1.5), Inches(12.33), [
    ["Scenario","Exit EV/EBITDA","EV (Rs cr)","Equity (Rs cr)","Equity (US$ bn)","Driver"],
    ["Bear","13x","~36,400","~37,900","~4.5","Prolonged AI de-rating, sub-7% growth"],
    ["Base","16x","~44,800","~46,300","~5.4","Growth recovery ~9–10%, margin 18%"],
    ["Bull","20x","~56,000","~57,500","~6.8","AI-services re-rating, double-digit growth"],
], col_widths=[Inches(1.5), Inches(2.0), Inches(1.9), Inches(2.0), Inches(2.0), Inches(2.93)], font=11.5, h=Inches(2.0))
txt(s, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.5),
    "Cross-checks: EV/Revenue ~3.0x (Base) and a DCF (WACC ~11.5%, g 4.5%) both support the ~Rs 44,000–50,000 cr Base equity range.",
    12, DGREY, italic=True)
txt(s, Inches(0.5), Inches(4.5), Inches(12), Inches(0.4), "Sponsor returns (Carlyle) — true multi-flow IRR", 14, NAVY, bold=True)
table(s, Inches(0.5), Inches(4.95), Inches(12.33), [
    ["Scenario","Realisations (IPO OFS + residual)","Hold","MOIC","Gross IRR"],
    ["Bear","Rs 8,750 cr + residual","~7 yrs","~1.7x","~8.6%"],
    ["Base","Rs 8,750 cr + ~Rs 34,300 cr","~7 yrs","~2.0x","~11.2%"],
    ["Bull","Rs 8,750 cr + ~Rs 42,600 cr","~7 yrs","~2.3x","~14.2%"],
], col_widths=[Inches(1.6), Inches(5.13), Inches(1.6), Inches(2.0), Inches(2.0)], font=11.5, h=Inches(1.7))

# A.9/10 Timing + recommendation
s = add_slide()
header(s, "HEXAWARE — SECTIONS 9–10", "Timing & Final Recommendation")
bullets(s, Inches(0.5), Inches(1.4), Inches(6.1), Inches(4.5), [
    "Industry mid-down-cycle in 2026: near-trough multiples (~15–18x P/E) on AI fears.",
    "Don't force a full exit into the trough — stagger on strength.",
    "Catalyst: growth re-acceleration to double digits would unlock the Bull case.",
    "Constructive capital markets + global rate easing support exit multiples into 2027–28.",
    "Recommended hold: ~7 yrs (2021–2028); staggered monetisation CY2026–CY2028.",
], size=13.5)
rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.2), LGOLD)
txt(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.4), "IC RECOMMENDATION", 14, NAVY, bold=True)
bullets(s, Inches(7.0), Inches(1.95), Inches(5.7), Inches(4.5), [
    "Verdict: HOLD & STAGGER-EXIT by ~CY2028.",
    "Route: staggered listed-market blocks; sponsor-to-sponsor as alternative.",
    "Buyers: FIIs/DIIs/SWFs (blocks); Blackstone/EQT/KKR; Cognizant/Capgemini (strategic).",
    "Base equity ~Rs 46,000 cr (~US$5.4bn).",
    "Expected ~2.0x MOIC / ~11% IRR (Base).",
    "Success factors: re-accelerate growth, expand margin to ~19–20%, manage overhang.",
], size=12.5, color=NAVY)

# =====================================================================
# PART B — SIEMENS
# =====================================================================
divider("PART B", "Siemens Ltd. (India)",
        "NSE: SIEMENS  •  Strategic parent: Siemens AG ~69%  •  framed as a HYPOTHETICAL PE study")

# B.1 Overview
s = add_slide()
header(s, "SIEMENS — SECTION 1", "Investment Overview (with control caveat)")
rect(s, Inches(0.5), Inches(1.3), Inches(12.33), Inches(0.75), RED)
txt(s, Inches(0.7), Inches(1.38), Inches(12), Inches(0.6),
    "CAVEAT: Siemens AG owns ~69%. A PE control transaction is effectively unavailable — analysis is a hypothetical PE / take-private study.",
    13, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
bullets(s, Inches(0.5), Inches(2.3), Inches(6.1), Inches(4.3), [
    "India's flagship industrial-tech & electrification leader (since 1922).",
    "Post Apr-2025 Energy demerger: continuing co. = Digital Industries, Smart Infrastructure, Mobility, Low-Voltage.",
    "Order backlog ~Rs 43,004 cr (Q1FY26) = strong revenue visibility.",
    "Net-cash; ~600% dividend (FY24); access to Siemens AG tech (Xcelerator, NVIDIA).",
], size=13)
bullets(s, Inches(6.8), Inches(2.3), Inches(6.0), Inches(4.3), [
    "Thesis: India capex/electrification supercycle (power demand >7% CAGR).",
    "Risks: controlling-shareholder lock, capex cyclicality, commodity/FX margin swings, premium valuation.",
    "Strategic relevance: a 'watch & learn' asset — PE access only via minority/PIPE blocks or future carve-outs.",
], size=13)

# B.2 Industry + Porter + SWOT
s = add_slide()
header(s, "SIEMENS — SECTION 2", "Industry, Porter's Five Forces & SWOT")
txt(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
    "Tailwinds: electrification, grid modernisation, 'Make in India' capex, metro/rail, data-centre power. Power demand >7% CAGR (Kroll).",
    11.5, DGREY, italic=True)
table(s, Inches(0.5), Inches(1.8), Inches(6.0), [
    ["Porter's Force","Intensity"],
    ["Competitive rivalry","High"],
    ["New entrants","Low"],
    ["Supplier power","Medium"],
    ["Buyer power","Medium-High"],
    ["Substitutes","Low"],
], col_widths=[Inches(4.0), Inches(2.0)], font=12, h=Inches(3.0))
txt(s, Inches(6.8), Inches(1.8), Inches(6), Inches(0.35), "SWOT", 14, NAVY, bold=True)
table(s, Inches(6.8), Inches(2.15), Inches(6.0), [
    ["Strengths","Weaknesses"],
    ["Market leadership, brand, backlog, net-cash","Cyclicality, commodity/FX, premium valuation"],
    ["Opportunities","Threats"],
    ["Grid/data-centre/rail capex, software annuities","Capex slowdown, de-rating, bidding pressure"],
], col_widths=[Inches(3.0), Inches(3.0)], font=10.5, first_col_label=False, h=Inches(2.6))
txt(s, Inches(0.5), Inches(5.6), Inches(6.1), Inches(1.4),
    "Net: a structurally attractive, HIGH-BARRIER oligopoly — superior to IT on entrant/substitution risk, but cyclical and richly valued.",
    12.5, NAVY, bold=True)

# B.3 Financials
s = add_slide()
header(s, "SIEMENS — SECTION 3", "5-Year Financials (FY Oct–Sep)")
txt(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.35),
    "FY20–FY24 include Energy (pre-demerger); not comparable to continuing-ops quarters below.", 11, DGREY, italic=True)
table(s, Inches(0.5), Inches(1.65), Inches(12.33), [
    ["Metric (Rs cr)","FY20","FY21","FY22","FY23","FY24"],
    ["Revenue","9,824","13,136","16,047","19,472","22,157"],
    ["Revenue growth %","—","+33.7%","+22.2%","+21.3%","+13.8%"],
    ["PAT","—","—","—","—","2,718"],
], col_widths=[Inches(2.93)] + [Inches(1.88)]*5, font=11.5, h=Inches(1.7))
txt(s, Inches(0.5), Inches(3.7), Inches(12), Inches(0.35), "Continuing operations (ex-Energy), recent quarters", 13, NAVY, bold=True)
table(s, Inches(0.5), Inches(4.1), Inches(9.0), [
    ["Period","Revenue (Rs cr)","Growth","PAT (Rs cr)"],
    ["Q4 FY25 (Sep-25)","5,171","+16%","485"],
    ["Q1 FY26 (Dec-25)","3,831","+14%","269"],
    ["Q2 FY26 (Mar-26)","n/d","+14.6%","355"],
], col_widths=[Inches(3.0), Inches(2.2), Inches(1.8), Inches(2.0)], font=11.5, h=Inches(1.7))
txt(s, Inches(0.5), Inches(6.6), Inches(12), Inches(0.4),
    "~23% pre-demerger revenue CAGR; continuing co. sustains mid-teens growth but VOLATILE PAT (commodities, FX, New Labour Codes). Sources: WSJ, Quartr, Siemens IR.",
    10.5, DGREY, italic=True)

# B.4 Value creation
s = add_slide()
header(s, "SIEMENS — SECTION 4", "Value Creation Opportunities (hypothetical lens)")
table(s, Inches(0.5), Inches(1.5), Inches(12.33), [
    ["Lever","Action / Impact"],
    ["Revenue expansion","Convert Rs 43,000 cr backlog; grid/data-centre/rail wins -> sustains mid-teens growth"],
    ["Margin improvement","Commodity hedging, pricing discipline, software/services mix -> +100–200 bps"],
    ["Digital transformation","Scale Xcelerator / industrial-AI annuities -> higher-multiple recurring revenue"],
    ["Capital structure","(Hypothetical) leverage — BLOCKED by strategic control"],
], col_widths=[Inches(3.2), Inches(9.13)], font=12, h=Inches(2.5))
rect(s, Inches(0.5), Inches(4.4), Inches(12.33), Inches(1.8), LGOLD)
txt(s, Inches(0.7), Inches(4.55), Inches(12), Inches(1.5),
    "The most powerful value lever — the SOTP unlock — was ALREADY executed by the parent via the 2025 Energy demerger "
    "(Siemens Energy India listed Jun-2025 at ~Rs 1.01–1.06 lakh cr; Jefferies >US$10bn). Limited residual PE 'alpha' remains.",
    14, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# B.5 Exit readiness
s = add_slide()
header(s, "SIEMENS — SECTION 5", "Exit Readiness — 9.0 / 10 (but no sponsor control)")
table(s, Inches(0.5), Inches(1.5), Inches(7.4), [
    ["Dimension","Score /10"],
    ["Management quality","9"],["Corporate governance","9"],
    ["Scalability","8"],["Market positioning","9"],
    ["Investor appeal","9"],["Strategic attractiveness","9"],
    ["Financial maturity","9"],["Public-market readiness","10"],
    ["OVERALL","9.0"],
], col_widths=[Inches(5.4), Inches(2.0)], font=11.5, h=Inches(4.6))
rect(s, Inches(8.2), Inches(1.5), Inches(4.6), Inches(4.6), GREY)
txt(s, Inches(8.4), Inches(1.7), Inches(4.2), Inches(4.2),
    "As a STANDALONE asset, Siemens India scores 9.0/10.\n\n"
    "But the binding constraint is ACCESS / CONTROL, not readiness.\n\n"
    "Siemens AG's ~69% stake means a PE-led exit is effectively unavailable.",
    15, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# B.6/7 Buyers + routes
s = add_slide()
header(s, "SIEMENS — SECTIONS 6–7", "Buyer Universe & Exit Routes")
txt(s, Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
    "Strategic buyers (ABB, Schneider, Hitachi, Mitsubishi, Eaton): ALL low probability — parent won't sell + antitrust hurdles.", 12, DGREY, italic=True)
table(s, Inches(0.5), Inches(1.8), Inches(12.33), [
    ["Exit route","Feasibility","P(success)"],
    ["Demerger-led SOTP value unlock (EXECUTED 2025)","High","Done"],
    ["Partial / staggered block sale (minority holders only)","Medium-High","~70%"],
    ["Strategic / control sale","Very Low","~5%"],
    ["Secondary buyout / sponsor-to-sponsor","Very Low","~5%"],
    ["Merger","Very Low","~5%"],
], col_widths=[Inches(7.5), Inches(2.83), Inches(2.0)], font=11.5, h=Inches(3.0))
rect(s, Inches(0.5), Inches(5.2), Inches(12.33), Inches(1.2), LBLUE)
txt(s, Inches(0.7), Inches(5.32), Inches(12), Inches(1.0),
    "The actionable 'route' that worked was the 2025 demerger (SOTP). For a non-parent minority holder, staggered block sales on strength are the only realistic monetisation.",
    13, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# B.8 Valuation
s = add_slide()
header(s, "SIEMENS — SECTION 8", "Exit Valuation & Hypothetical Returns")
txt(s, Inches(0.5), Inches(1.25), Inches(12), Inches(0.4),
    "Current continuing-co. market cap ~Rs 1.27–1.33 lakh cr (~US$15bn) at a steep ~60–70x P/E. We model exit-year (GROWN) PAT over ~4 yrs.", 11.5, DGREY, italic=True)
table(s, Inches(0.5), Inches(1.75), Inches(12.33), [
    ["Scenario","Exit PAT (Rs cr)","Exit P/E","Equity (Rs cr)","Equity (US$ bn)","MOIC","IRR"],
    ["Bear","~3,000","32x","~96,000","~11.3","~0.74x","~ -7%"],
    ["Base","~3,700","45x","~166,500","~19.6","~1.28x","~6.4%"],
    ["Bull","~4,400","52x","~228,800","~26.9","~1.76x","~15.2%"],
], col_widths=[Inches(1.5), Inches(1.9), Inches(1.5), Inches(2.1), Inches(2.1), Inches(1.6), Inches(1.63)], font=11, h=Inches(2.0))
rect(s, Inches(0.5), Inches(4.1), Inches(12.33), Inches(2.2), LGOLD)
txt(s, Inches(0.7), Inches(4.25), Inches(12), Inches(1.9),
    "Key insight: the BEAR case (~Rs 96,000 cr) sits BELOW today's ~Rs 1.30 lakh cr market cap — a financial buyer entering at the current "
    "premium would LOSE money if multiples de-rate. Returns are modest-to-poor because the asset is already richly valued and the SOTP "
    "alpha has been captured. Hold assumed 4 yrs; entry at current mcap (Rs 130,000 cr). Figures computed live in the Excel model.",
    13.5, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

# B.9/10 Timing + recommendation
s = add_slide()
header(s, "SIEMENS — SECTIONS 9–10", "Timing & Final Recommendation")
bullets(s, Inches(0.5), Inches(1.4), Inches(6.1), Inches(4.5), [
    "India in an electrification/capex up-cycle — supportive, but valuations ELEVATED.",
    "Capital-goods mcap pulled back ~14% in late CY25 (Kroll) — de-rating risk.",
    "Value-creation event (demerger) already realised in 2025.",
    "Parent (Siemens AG): retain and compound — no exit recommendation.",
    "Minority holder: monetise opportunistically on strength within 2–4 years.",
], size=13.5)
rect(s, Inches(6.8), Inches(1.4), Inches(6.0), Inches(5.2), RED)
txt(s, Inches(7.0), Inches(1.5), Inches(5.6), Inches(0.4), "IC RECOMMENDATION", 14, WHITE, bold=True)
bullets(s, Inches(7.0), Inches(1.95), Inches(5.7), Inches(4.5), [
    "Verdict: DO NOT PURSUE A CONTROL INVESTMENT — PASS.",
    "Structurally unavailable (Siemens AG ~69%) and already richly valued.",
    "Base equity ~Rs 166,500 cr (~US$19.6bn); hypothetical ~1.28x / ~6.4%.",
    "Downside real: Bear MOIC ~0.74x (value loss).",
    "If exposure desired: minority/PIPE blocks only, monetised near-term.",
    "Lasting PE lesson: the 2025 demerger as an SOTP value-unlock template.",
], size=12.5, color=WHITE)

# =====================================================================
# CLOSING — comparative + final
# =====================================================================
s = add_slide()
header(s, "PART C", "Comparative Scorecard")
table(s, Inches(0.5), Inches(1.5), Inches(12.33), [
    ["Metric","Hexaware","Siemens India"],
    ["5-yr revenue CAGR","~17% (CY21–25)","~23% (FY20–24, pre-demerger)"],
    ["Margin profile","~16–17% EBITDA (expansion lever)","~11–13% EBIT (cyclical)"],
    ["Exit readiness","8.4 / 10","9.0 / 10 (no sponsor control)"],
    ["Control available to PE?","YES","NO"],
    ["Base equity value","~Rs 46,000 cr (~US$5.4bn)","~Rs 166,500 cr (~US$19.6bn)"],
    ["Base MOIC / IRR","~2.0x / ~11%","~1.28x / ~6.4% (hypothetical)"],
    ["Optimal route","Staggered blocks -> full exit ~CY2028","Demerger SOTP (done); blocks only"],
    ["IC VERDICT","HOLD & STAGGER-EXIT","PASS on control"],
], col_widths=[Inches(3.0), Inches(4.66), Inches(4.66)], font=12, h=Inches(4.7))

s = add_slide()
rect(s, 0, 0, SW, SH, NAVY)
rect(s, Inches(0.8), Inches(1.4), Inches(2.2), Inches(0.10), GOLD)
txt(s, Inches(0.8), Inches(0.7), Inches(11), Inches(0.5), "FINAL IC RECOMMENDATION", 16, GOLD, bold=True)
txt(s, Inches(0.8), Inches(1.7), Inches(11.8), Inches(1.0),
    "Hexaware: HOLD & STAGGER-EXIT by ~CY2028", 28, WHITE, bold=True)
txt(s, Inches(0.8), Inches(2.6), Inches(11.8), Inches(1.0),
    "Avoid the 2026 trough; stagger listed-market blocks to a Base ~2.0x MOIC / ~11% IRR; "
    "retain sponsor-to-sponsor & strategic-sale optionality.", 15, LBLUE)
txt(s, Inches(0.8), Inches(3.9), Inches(11.8), Inches(1.0),
    "Siemens India: PASS on control", 28, WHITE, bold=True)
txt(s, Inches(0.8), Inches(4.8), Inches(11.8), Inches(1.2),
    "A best-in-class, ~9/10 exit-ready asset, but structurally unavailable to sponsors (Siemens AG ~69%) and already richly "
    "valued (Bear case below today's price). Study the 2025 demerger as the value-unlock template.", 15, LBLUE)
rect(s, Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.7), GOLD)
txt(s, Inches(1.0), Inches(6.38), Inches(11.3), Inches(0.55),
    "Companion model: PE_Exit_Model_Hexaware_Siemens.xlsx — Assumptions feed Valuation & Returns; live IC Dashboard.",
    13, NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)

prs.save("/projects/sandbox/PE_Exit_Strategy_Hexaware_Siemens.pptx")
print("PPTX saved with", len(prs.slides._sldIdLst), "slides")
